"""Utility functions: seeding, logging, file I/O helpers."""

import random
import sys
import numpy as np
import torch
import os
import json


def set_seed(seed: int = 42):
    """Set all random seeds for reproducibility."""
    random.seed(seed)
    np.random.seed(seed)
    torch.manual_seed(seed)
    if torch.cuda.is_available():
        torch.cuda.manual_seed(seed)
        torch.cuda.manual_seed_all(seed)
    # For transformers
    os.environ["PYTHONHASHSEED"] = str(seed)


def setup_logging(verbose: bool = True):
    """Configure basic logging for the project."""
    import logging
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(
        level=level,
        format="%(asctime)s [%(levelname)s] %(message)s",
        datefmt="%H:%M:%S",
    )


def save_json(data, path: str):
    """Save data to JSON file, creating directories if needed."""
    os.makedirs(os.path.dirname(path) if os.path.dirname(path) else ".", exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def load_json(path: str):
    """Load data from JSON file."""
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


# Force UTF-8 for stdout/stderr on Windows
if hasattr(sys, "stdout"):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass


def format_duration(seconds: float) -> str:
    """Format a duration in seconds to a human-readable string."""
    if seconds < 60:
        return f"{seconds:.1f}s"
    minutes = int(seconds // 60)
    secs = seconds % 60
    return f"{minutes}m {secs:.0f}s"


# ── File parser (for web upload) ──────────────────────────

def parse_file(filepath: str):
    """Auto-detect format and extract text. Returns None if unsupported."""
    ext = os.path.splitext(filepath)[1].lower()
    parsers = {
        ".txt": lambda p: _read(p), ".md": lambda p: _read(p),
        ".csv": _parse_csv, ".json": _parse_json, ".xml": _read,
        ".pdf": _parse_pdf, ".pptx": _parse_pptx, ".docx": _parse_docx,
    }
    parser = parsers.get(ext)
    if not parser:
        return None
    try:
        return parser(filepath)
    except Exception as e:
        return f"[Error reading {ext}: {e}]"


def _read(path):
    with open(path, "r", encoding="utf-8", errors="surrogateescape") as f:
        return f.read()


def _parse_csv(path):
    import csv, io
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return "\n".join(" | ".join(row) for row in csv.reader(f))


def _parse_json(path):
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    if isinstance(data, (dict, list)):
        if isinstance(data, dict) and ("dialogue" in data or "text" in data):
            return data.get("dialogue") or data.get("text", "")
        return json.dumps(data, ensure_ascii=False, indent=2)
    return str(data)


def _parse_pdf(path):
    from PyPDF2 import PdfReader
    return "\n\n".join(p.extract_text() or "" for p in PdfReader(path).pages)


def _parse_pptx(path):
    from pptx import Presentation
    slides = []
    for i, slide in enumerate(Presentation(path).slides, 1):
        texts = [t for shape in slide.shapes if shape.has_text_frame
                 for para in shape.text_frame.paragraphs if (t := para.text.strip())]
        if texts:
            slides.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _parse_docx(path):
    from docx import Document
    return "\n\n".join(p.text for p in Document(path).paragraphs if p.text.strip())
